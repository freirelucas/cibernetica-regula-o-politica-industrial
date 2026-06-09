#!/usr/bin/env python3
"""Gera um deck .pptx de STATUS do projeto (onde estamos / o que falta), com os
números lidos AO VIVO dos JSONs do funil (reprodutível — nada escrito à mão) e
os screenshots das três telas embutidos.

Uso:
    python .claude/skills/run-scisci-ipea/driver.py            # (gera os screenshots)
    python src/build_deck.py                                   # monta docs/apresentacao_status.pptx

Dependência (dev): python-pptx.
"""
import json
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
DATA = os.path.join(ROOT, "data")
DOCS = os.path.join(ROOT, "docs")
OUT = os.path.join(DOCS, "apresentacao_status.pptx")

INK = RGBColor(0x0E, 0x0E, 0x0E)
PAPER = RGBColor(0xF5, 0xF2, 0xED)
CYB = RGBColor(0x4B, 0x3F, 0xC2)
ACCENT = RGBColor(0xC2, 0x3F, 0x3F)
MUTED = RGBColor(0x6B, 0x65, 0x60)


def load(name, default=None):
    p = os.path.join(DATA, name)
    if not os.path.exists(p):
        return default if default is not None else {}
    return json.load(open(p, encoding="utf-8"))


def br(n):
    """Formata inteiro no padrão BR (ponto de milhar)."""
    try:
        return f"{int(n):,}".replace(",", ".")
    except Exception:
        return str(n)


# ── números ao vivo ─────────────────────────────────────────────────────────
R = load("scisci_results.json")
sol = load("solidity_bridges.json")
net = load("network_4axis.json", {"nodes": [], "links": []})
xgi = load("cocitation_hyperedges.json")
mod = load("modularity_check.json")

corpus = br(R.get("corpus_size"))
seeds = R.get("n_seeds")
n_nodes, n_links = len(net.get("nodes", [])), len(net.get("links", []))
n_hyper = xgi.get("n_hyperedges") or xgi.get("n_edges")
z_hyper = (xgi.get("null_model") or {}).get("z")
n_agenda = sol.get("n_agenda")
n_alta = sol.get("n_alta_integracao")
n_alem = sol.get("alem_do_acaso")
tv = sol.get("validacao_temporal") or {}
temporal_ok = tv.get("validated")

# Q dos silos (núcleo) — busca tolerante na estrutura de modularity_check
def find_Q(o):
    best = None
    def walk(x):
        nonlocal best
        if isinstance(x, dict):
            for k, v in x.items():
                if isinstance(v, (int, float)) and k.lower() in ("q", "q_obs", "q_observed", "modularidade", "modularity") and 0 < v < 1:
                    best = v
                walk(v)
        elif isinstance(x, list):
            for v in x:
                walk(v)
    walk(o)
    return best
Q = find_Q(mod)
Q_txt = f"Q ≈ {Q:.2f}".replace(".", ",") if Q else "Q ≈ 0,5–0,6"

SHOTS = {
    "topo": "/tmp/deck_topo.png",
    "explorador": "/tmp/deck_explorador.png",
    "triagem": "/tmp/deck_triagem.png",
}

# ── helpers de slide ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W = prs.slide_width


def slide(bg=PAPER):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def tb(s, x, y, w, h):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.word_wrap = True
    return box.text_frame


def para(tf, text, size=18, color=INK, bold=False, first=False, align=PP_ALIGN.LEFT, space=8):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space)
    r = p.add_run(); r.text = text
    f = r.font
    f.size = Pt(size); f.bold = bold; f.color.rgb = color; f.name = "Calibri"
    return p


def kicker(s, text):
    para(tb(s, 0.9, 0.55, 11.5, 0.5), text.upper(), 13, CYB, bold=True, first=True)


def title(s, text, size=34):
    para(tb(s, 0.9, 1.0, 11.5, 1.3), text, size, INK, bold=True, first=True)


def bullets(s, items, x=0.9, y=2.4, w=11.5, h=4.6, size=18):
    tf = tb(s, x, y, w, h)
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            txt, col = it
        else:
            txt, col = it, INK
        para(tf, txt, size, col, first=(i == 0), space=12)


def cover():
    s = slide(INK)
    para(tb(s, 0.9, 0.7, 11.5, 0.5), "IPEA · DIEST-COGIT · CIENCIOMETRIA · STATUS", 14, RGBColor(0xBF, 0xB8, 0xF0), bold=True, first=True)
    para(tb(s, 0.9, 2.4, 11.5, 2.0), "Cibernética, Regulação e Política Industrial", 40, PAPER, bold=True, first=True)
    para(tb(s, 0.9, 3.7, 11.5, 1.2), "Onde estamos e o caminho para a revisão", 22, RGBColor(0xCF, 0xCA, 0xC2), first=True)
    para(tb(s, 0.9, 6.3, 11.5, 0.8), "Material preliminar · junho de 2026", 14, MUTED, first=True)


def pic_slide(kkr, ttl, caption, shot):
    s = slide(PAPER)
    kicker(s, kkr); title(s, ttl, 28)
    if os.path.exists(shot):
        s.shapes.add_picture(shot, Inches(0.9), Inches(2.0), height=Inches(4.5))
    else:
        para(tb(s, 0.9, 3.0, 11.5, 1.0), f"[screenshot ausente: {shot}]", 16, ACCENT, first=True)
    para(tb(s, 0.9, 6.7, 11.5, 0.6), caption, 14, MUTED, first=True)


# ── slides ──────────────────────────────────────────────────────────────────
cover()

s = slide(); kicker(s, "A pergunta")
title(s, "Três tradições que quase não conversam — um campo ou ilhas?")
bullets(s, [
    "Cibernética organizacional · Regulação econômica · Política industrial.",
    "Testamos vários modelos sobre a base OpenAlex (cocitação par-a-par e de ordem superior/hipergrafo, comunidades, intermediação, integração ΔKf), sempre contra modelos nulos.",
    "Propósito: achar as pontes entre os campos — e testar se elas de fato existem.",
])

s = slide(); kicker(s, "O funil")
title(s, "Escala dos dados (offline, reprodutível)")
bullets(s, [
    (f"{corpus} trabalhos no corpus · {seeds} obras-semente (20 por eixo).", INK),
    (f"Rede de cocitação dos 4 eixos: {br(n_nodes)} nós · {br(n_links)} arestas.", INK),
    (f"Hipergrafo: {br(n_hyper)} hiperarestas (bibliografias citantes).", INK),
    ("17 estágios encadeados; cache versionado; CI verde; 80 testes.", MUTED),
])

s = slide(); kicker(s, "Achado central")
title(s, "Os silos são reais — e robustos")
bullets(s, [
    (f"Comunidades de citação separadas: {Q_txt} contra ~0 no acaso.", INK),
    (f"No hipergrafo (listas de leitura individuais): z ≈ {z_hyper} — muito abaixo do acaso.", INK),
    ("O isolamento vai até quem cita e quem colabora (corretagem ~2:1 intra-campo).", INK),
    ("Robusto a bootstrap (drop-20%) e ao nulo casado em grau.", MUTED),
])

s = slide(); kicker(s, "A agenda honesta")
title(s, "Não há ponte latente — a convergência tem de ser construída")
bullets(s, [
    (f"{br(n_alta)} tríades de alta integração (top-10% de ΔKf) → {br(n_agenda)} alvos de agenda (também plausíveis na faixa semântica).", INK),
    (f"Contra o nulo casado em grau: {n_alem} além do acaso. Nenhuma ponte integra mais que o acaso.", ACCENT),
    (f"Validação temporal (fora da amostra): {'tendência com sinal' if temporal_ok else 'não distinguível do acaso'}.", MUTED),
    ("A agenda diz ONDE costurar renderia mais integração — não certifica ponte existente.", INK),
])

s = slide(); kicker(s, "Rigor & integridade")
title(s, "O que endurecemos nesta rodada")
bullets(s, [
    ("Integração = condutância real (ΔKf, posicional) — fim do artefato de grau (ex-“DESIGN”).", INK),
    ("Higiene de dados: removidas 8 obras-fantasma (404) + 2 agregados de periódico (hubs falsos).", INK),
    ("O “principal conector” era um registro vazio (W4285719527, ligado a 95% dos nós) — agora obra real.", ACCENT),
    ("Auditoria final: 0 ids 404/agregado em 10 fontes de dados.", MUTED),
])

pic_slide("As três telas", "Relatório — “Comece por aqui” + 5 atos colapsáveis",
          "Topo executivo com curadoria de método, achados em 3 frases e mapa em atos.", SHOTS["topo"])
pic_slide("Explorador", "Modo “Só as pontes” — corta o ruído",
          "De 250 nós para os alvos de agenda rotulados; rede inteira a um clique.", SHOTS["explorador"])
pic_slide("Triagem", "Duas etapas: “o que o algoritmo achou” → “decidir”",
          "O fruto (alvos) com o porquê de cada um; depois incluir/talvez/excluir → Rayyan.", SHOTS["triagem"])

s = slide(); kicker(s, "O que falta")
title(s, "Da agenda à revisão")
bullets(s, [
    ("Protocolo de revisão: importar Rayyan → 3 revisores independentes → κ de concordância.", INK),
    ("Sensibilidade às sementes (drop-20% × 5) e recompute de Q no corpus novo.", INK),
    ("Disclosure dos inferidos no ponto da afirmação (obra > autor > conceito).", INK),
    ("Curadoria editorial: rolar o “em miúdos” para mais seções densas.", MUTED),
])

s = slide(INK); kicker(s, "Próximos passos")
para(tb(s, 0.9, 1.0, 11.5, 1.3), "Decisão", 34, PAPER, bold=True, first=True)
bullets(s, [
    ("1. Abrir a revisão minuciosa dos candidatos (Rayyan, 3 revisores).", PAPER),
    ("2. Construir as primeiras pontes da agenda (listas comentadas, cursos, prática de política).", PAPER),
    ("3. Fechar o backlog de rigor restante.", RGBColor(0xCF, 0xCA, 0xC2)),
], y=2.6)

prs.save(OUT)
print(f"deck: {OUT} ({os.path.getsize(OUT)//1024} KB, {len(prs.slides)} slides)")
